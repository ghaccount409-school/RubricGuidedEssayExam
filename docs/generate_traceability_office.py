"""
Generate Word (.docx) and PowerPoint (.pptx) from traceability content.
Run from repo root: python docs/generate_traceability_office.py
Outputs: docs/requirements-test-traceability.docx, docs/requirements-test-traceability.pptx
"""
from __future__ import annotations

import importlib.util
from pathlib import Path

from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt

DOCS = Path(__file__).resolve().parent


def _load_traceability():
    path = DOCS / "traceability_data.py"
    spec = importlib.util.spec_from_file_location("traceability_data", path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Cannot load {path}")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod.ROWS_MAIN, mod.ROWS_UNIT


ROWS_MAIN, ROWS_UNIT = _load_traceability()


def build_docx() -> None:
    doc = Document()
    sec = doc.sections[0]
    sec.orientation = WD_ORIENT.LANDSCAPE
    sec.page_width, sec.page_height = sec.page_height, sec.page_width
    t = doc.add_heading("Requirements-to-Tests Traceability", 0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p = doc.add_paragraph()
    p.add_run("Project: ").bold = True
    p.add_run("RubricGuidedEssayExam")
    doc.add_paragraph(
        "Repository: https://github.com/ALGeek01/RubricGuidedEssayExam"
    )
    doc.add_paragraph(
        "SRS reference: requirements-specification-csc394-capstone.md (v1.1, local)."
    )
    doc.add_paragraph(
        "Tests use MOCK_LLM=1 and isolated SQLite (see tests/conftest.py)."
    )

    doc.add_heading("1. Test inventory", level=1)
    inv = doc.add_table(rows=5, cols=2)
    inv.style = "Table Grid"
    inv_data = [
        ("tests/conftest.py", "SQLite, MOCK_LLM, client fixture, DB reset"),
        ("tests/general/test_unit.py", "pytest.mark.unit"),
        ("tests/general/test_api.py", "pytest.mark.integration"),
        ("tests/security/test_http.py", "pytest.mark.security"),
    ]
    inv.rows[0].cells[0].text = "File"
    inv.rows[0].cells[1].text = "Role"
    for i, (a, b) in enumerate(inv_data, start=1):
        inv.rows[i].cells[0].text = a
        inv.rows[i].cells[1].text = b

    doc.add_heading("2. SRS → tests (main mapping, with explanations)", level=1)
    tbl = doc.add_table(rows=1 + len(ROWS_MAIN), cols=4)
    tbl.style = "Table Grid"
    hdr = tbl.rows[0].cells
    hdr[0].text = "Requirement"
    hdr[1].text = "Test(s)"
    hdr[2].text = "Coverage note"
    hdr[3].text = "Explanation"
    for r, row in enumerate(ROWS_MAIN, start=1):
        tbl.rows[r].cells[0].text = row[0]
        tbl.rows[r].cells[1].text = row[1]
        tbl.rows[r].cells[2].text = row[2]
        tbl.rows[r].cells[3].text = row[3]
    for cell in hdr:
        for para in cell.paragraphs:
            for run in para.runs:
                run.bold = True

    doc.add_heading("3. Unit tests vs SRS (with explanations)", level=1)
    ut = doc.add_table(rows=1 + len(ROWS_UNIT), cols=3)
    ut.style = "Table Grid"
    ut.rows[0].cells[0].text = "Test"
    ut.rows[0].cells[1].text = "SRS relationship"
    ut.rows[0].cells[2].text = "Explanation"
    for r, row in enumerate(ROWS_UNIT, start=1):
        ut.rows[r].cells[0].text = row[0]
        ut.rows[r].cells[1].text = row[1]
        ut.rows[r].cells[2].text = row[2]
    for c in range(3):
        for para in ut.rows[0].cells[c].paragraphs:
            for run in para.runs:
                run.bold = True

    doc.add_heading("4. Security tests", level=1)
    doc.add_paragraph(
        "test_not_found_exam_question_has_no_python_traceback; "
        "test_not_found_exam_results_has_no_python_traceback; "
        "test_exam_start_rejects_empty_student_id_without_traceback; "
        "test_professor_unknown_exam_is_404 — NFR-LOG-1 partial (no traceback in HTML)."
    )

    doc.add_heading("5. Summary", level=1)
    doc.add_paragraph(
        "Well covered for current scope: basic exam flow, two questions, 404s, "
        "professor smoke tests, no traceback on errors, duplicate answer rejected."
    )
    doc.add_paragraph(
        "Major gaps: FR-GEN-1 sandbox/exec, FR-SEC-*, full FR-STU-*, schema grading, "
        "follow-ups, DB retention, final grade schemes, disputes, export, most NFRs."
    )

    out = DOCS / "requirements-test-traceability.docx"
    doc.save(out)
    print(f"Wrote {out}")


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    sub = slide.placeholders[1]
    sub.text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = PptxPt(18)


def build_pptx() -> None:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    add_title_slide(
        prs,
        "Requirements ↔ Test Traceability",
        "RubricGuidedEssayExam · SRS v1.1 · 2026-04-07",
    )

    add_bullet_slide(
        prs,
        "Purpose",
        [
            "Map capstone SRS requirements to pytest tests",
            "SRS: requirements-specification-csc394-capstone.md",
            "Tests: github.com/ALGeek01/RubricGuidedEssayExam/tree/main/tests",
            "MOCK_LLM=1 — no live Together API in tests",
        ],
    )

    add_bullet_slide(
        prs,
        "Test files",
        [
            "conftest.py — SQLite + DB reset + TestClient",
            "general/test_unit.py — education levels, TogetherApiError",
            "general/test_api.py — exam flows, professor pages",
            "security/test_http.py — no traceback in error HTML",
        ],
    )

    add_bullet_slide(
        prs,
        "Covered today (high level)",
        [
            "Exam start → question → answer → results (mock LLM)",
            "Two-question flow to results",
            "404 for missing exam / professor detail",
            "Professor dashboard + detail smoke (200)",
            "Reject second POST /answer after completion (400)",
            "Invalid education level / blank student id (400)",
            "404 responses without Python traceback in body",
        ],
    )

    add_bullet_slide(
        prs,
        "Gaps vs full SRS",
        [
            "FR-GEN-1: LLM Python + sandbox + schema validation",
            "FR-SEC-1 / FR-SEC-2: sandbox policy, authentication",
            "FR-STU-3/4/5: length limits, autosave, timers",
            "FR-GRADE-2/3, FR-FU-*, FR-DB-*, FR-FINAL-*, FR-DISP-2/3",
            "FR-REVIEW-3 export; NFR-PERF, NFR-REL, NFR-AUDIT, WCAG",
            "ADV-CURVE-1, ADV-DISC-1",
        ],
    )

    add_bullet_slide(
        prs,
        "Key test function names",
        [
            "test_full_exam_single_question_flow",
            "test_two_question_flow",
            "test_exam_start_redirects_and_creates_session",
            "test_completed_session_question_redirects_to_results",
            "test_answer_invalid_session",
            "test_professor_dashboard_and_detail",
            "test_not_found_exam_question_has_no_python_traceback",
        ],
    )

    add_bullet_slide(
        prs,
        "Deliverables",
        [
            "Markdown: docs/requirements-test-traceability.md",
            "Word: docs/requirements-test-traceability.docx",
            "PowerPoint: docs/requirements-test-traceability.pptx",
            "Regenerate Office files: python docs/generate_traceability_office.py",
        ],
    )

    out = DOCS / "requirements-test-traceability.pptx"
    prs.save(out)
    print(f"Wrote {out}")


def main() -> None:
    build_docx()
    build_pptx()


if __name__ == "__main__":
    main()
